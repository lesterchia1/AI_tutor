import streamlit as st
import numpy as np
import os
import groq
import uuid
import chardet
import fitz  # PyMuPDF
import docx
import gtts
from pptx import Presentation
import re
import json
import pickle

from langchain_groq import ChatGroq
from langchain_core.messages import HumanMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

# -----------------------------
# Init Models
# -----------------------------
st.set_page_config(page_title="AI Tutor", layout="wide")

# Initialize Groq client (for transcription + chat)
groq_client = groq.Groq(api_key=os.getenv("GROQ_API_KEY"))

chat_model = ChatGroq(model_name="llama-3.3-70b-versatile", api_key=os.getenv("GROQ_API_KEY"))

# Simple document storage instead of Chroma
os.makedirs("documents", exist_ok=True)
document_store = []

chat_memory = []

quiz_prompt = """
You are an AI assistant specialized in education and assessment creation. Create a comprehensive quiz based on the provided document content. Include multiple choice questions, true/false questions, and short answer questions. Format the quiz clearly with sections and provide answers at the end.
"""

# -----------------------------
# Helper Functions
# -----------------------------
def transcribe_audio(file_path):
    """Transcribe audio using Groq Whisper API"""
    with open(file_path, "rb") as f:
        transcript = groq_client.audio.transcriptions.create(
            file=f,
            model="whisper-large-v3"
        )
    return transcript.text

def clean_response(response):
    cleaned_text = re.sub(r"<think>.*?</think>", "", response, flags=re.DOTALL)
    cleaned_text = re.sub(r"(\*\*|\*|\[|\])", "", cleaned_text)
    cleaned_text = re.sub(r"^##+\s*", "", cleaned_text, flags=re.MULTILINE)
    cleaned_text = re.sub(r"\\", "", cleaned_text)
    cleaned_text = re.sub(r"---", "", cleaned_text)
    return cleaned_text.strip()

def generate_quiz(content):
    # Limit content length to avoid token limits
    if len(content) > 4000:
        content = content[:4000] + "..."
    
    prompt = f"{quiz_prompt}\n\nDocument content:\n{content}"
    response = chat_model.invoke(prompt)  # ← CHANGED TO invoke()
    return clean_response(response.content)

def retrieve_documents(query):
    """Simple keyword-based document retrieval"""
    if not document_store:
        return ["No documents available. Please upload documents first."]
    
    # Simple keyword matching (case insensitive)
    query_words = query.lower().split()
    relevant_docs = []
    
    for doc in document_store:
        doc_text = doc.lower()
        matches = sum(1 for word in query_words if word in doc_text)
        if matches > 0:
            relevant_docs.append(doc)
            if len(relevant_docs) >= 3:  # Limit to 3 documents
                break
    
    return relevant_docs if relevant_docs else ["No relevant documents found for your query."]

def chat_with_groq(user_input):
    relevant_docs = retrieve_documents(user_input)
    context = "\n".join(relevant_docs) if relevant_docs else "No relevant documents found."
    system_prompt = "You are a helpful AI assistant. Answer questions accurately and concisely."
    conversation_history = "\n".join(chat_memory[-10:]) if chat_memory else ""
    prompt = f"{system_prompt}\n\nConversation History:\n{conversation_history}\n\nUser Input: {user_input}\n\nContext:\n{context}"

    response = chat_model.invoke(prompt)
    cleaned_response = clean_response(response.content)

    chat_memory.append(f"User: {user_input}")
    chat_memory.append(f"AI: {cleaned_response}")

    return cleaned_response

def speech_playback(text):
    try:
        filename = f"output_{uuid.uuid4()}.mp3"
        tts = gtts.gTTS(text, lang="en")
        tts.save(filename)
        return filename
    except Exception as e:
        st.error(f"TTS Error: {e}")
        return None

def detect_encoding(file_path):
    with open(file_path, "rb") as f:
        raw_data = f.read(4096)
        detected = chardet.detect(raw_data)
        return detected["encoding"] or "utf-8"

def extract_text_from_pdf(path):
    try:
        doc = fitz.open(path)
        return "\n".join([page.get_text("text") for page in doc])
    except Exception as e:
        st.error(f"PDF extraction error: {e}")
        return ""

def extract_text_from_docx(path):
    try:
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except Exception as e:
        st.error(f"DOCX extraction error: {e}")
        return ""

def extract_text_from_pptx(path):
    try:
        presentation = Presentation(path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        st.error(f"PPTX extraction error: {e}")
        return ""

def process_document(file):
    try:
        ext = os.path.splitext(file.name)[-1].lower()
        content = ""
        
        if ext == ".pdf":
            content = extract_text_from_pdf(file.name)
        elif ext == ".docx":
            content = extract_text_from_docx(file.name)
        elif ext == ".pptx":
            content = extract_text_from_pptx(file.name)
        else:
            encoding = detect_encoding(file.name)
            with open(file.name, "r", encoding=encoding, errors="replace") as f:
                content = f.read()

        if not content.strip():
            return "Error: Could not extract text from the document."

        # Store document content for retrieval
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = text_splitter.split_text(content)
        
        # Add chunks to document store
        document_store.extend(chunks)
        
        # Save document store to file
        with open("documents/document_store.pkl", "wb") as f:
            pickle.dump(document_store, f)
        
        return generate_quiz(content)
        
    except Exception as e:
        return f"Error processing document: {str(e)}"

# Load existing documents if available
try:
    if os.path.exists("documents/document_store.pkl"):
        with open("documents/document_store.pkl", "rb") as f:
            document_store = pickle.load(f)
except:
    document_store = []

# -----------------------------
# Streamlit UI
# -----------------------------
st.title("📚 AI Tutor - Streamlit POC")

tab1, tab2, tab3, tab4 = st.tabs(["💬 AI Chatbot", "📄 Upload & Quiz", "🎥 Intro Video", "🎤 Audio Transcription"])

# --- Tab 1: Chatbot ---
with tab1:
    st.subheader("Chat with AI Tutor")
    user_input = st.text_input("Ask a question:")
    if st.button("Send") and user_input:
        with st.spinner("Thinking..."):
            reply = chat_with_groq(user_input)
        st.markdown(f"**AI:** {reply}")
        audio_file = speech_playback(reply)
        if audio_file:
            st.audio(audio_file)

# --- Tab 2: Upload & Quiz ---
with tab2:
    st.subheader("Upload Notes & Generate Quiz")
    uploaded_file = st.file_uploader("Upload a file (PDF, DOCX, PPTX, TXT)", type=["pdf", "docx", "pptx", "txt"])
    if uploaded_file:
        with st.spinner("Processing document and generating quiz..."):
            with open(uploaded_file.name, "wb") as f:
                f.write(uploaded_file.read())
            quiz = process_document(uploaded_file)
        st.text_area("Generated Quiz", quiz, height=400)
        
        # Show document stats
        st.info(f"Document processed. Total document chunks in memory: {len(document_store)}")

# --- Tab 3: Intro Video ---
with tab3:
    st.subheader("Introduction Video")
    st.video("https://github.com/lesterchia1/AI_tutor/raw/main/We%20not%20me%20video.mp4")

# --- Tab 4: Audio Transcription ---
with tab4:
    st.subheader("Upload an Audio File for Transcription")
    audio_file = st.file_uploader("Upload an audio file", type=["mp3", "wav", "m4a"])
    if audio_file:
        with st.spinner("Transcribing audio..."):
            with open(audio_file.name, "wb") as f:
                f.write(audio_file.read())
            transcript_text = transcribe_audio(audio_file.name)
        st.text_area("Transcription", transcript_text, height=300)
